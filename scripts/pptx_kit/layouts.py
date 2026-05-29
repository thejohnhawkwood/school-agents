"""
Reusable slide layouts — blank-slide compositions with consistent hierarchy.
All text is native PPTX; images pre-rendered via pptx_kit.images for correct aspect.
"""

from __future__ import annotations

from pathlib import Path
from typing import Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from pptx_kit.images import add_image_cover
from pptx_kit.text_fit import (
    fit_bullets_to_box,
    summarize_bullet,
    truncate_title,
)
from pptx_kit.tokens import DEFAULT_TOKENS, DesignTokens


def add_blank_slide(prs: Presentation):
    """Prefer true Blank layout when template provides it."""
    layouts = prs.slide_layouts
    for idx in (6, 5, 1):
        if idx < len(layouts):
            return prs.slides.add_slide(layouts[idx])
    return prs.slides.add_slide(layouts[0])


def apply_slide_background(slide, rgb, tokens: DesignTokens = DEFAULT_TOKENS) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def clear_tf(tf) -> None:
    try:
        tf.clear()
    except AttributeError:
        while len(tf.paragraphs) > 1:
            p = tf.paragraphs[-1]._element
            p.getparent().remove(p)
        if tf.paragraphs:
            tf.paragraphs[0].text = ""


def configure_text_frame(tf, tokens: DesignTokens = DEFAULT_TOKENS) -> None:
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)


def title_region_height_in(tokens: DesignTokens, title_pt: int, subtitle_pt: int | None) -> float:
    """Approximate vertical space title block consumes."""
    from pptx_kit.text_fit import pt_to_inches

    h = pt_to_inches(title_pt * tokens.line_spacing * 1.35)
    if subtitle_pt:
        h += pt_to_inches(subtitle_pt * tokens.line_spacing * 1.2)
    return h + 0.15


# --- 1. Title slide ---
def layout_title_slide(
    prs: Presentation,
    tokens: DesignTokens,
    tmp_dir: Path,
    temps: list[Path],
    *,
    title: str,
    subtitle: str,
    hero_image: Path | None = None,
    speaker_notes: str = "",
) -> None:
    slide = add_blank_slide(prs)
    apply_slide_background(slide, tokens.color_bg, tokens)
    tw = tokens.slide_width_in
    th = tokens.slide_height_in
    title = truncate_title(title, tokens.max_title_chars)

    if hero_image and hero_image.exists():
        _pic, tmp = add_image_cover(
            slide,
            hero_image,
            Inches(0),
            Inches(0),
            Inches(tw),
            Inches(th),
            tmp_dir,
        )
        temps.append(tmp)
        # Dark overlay band at bottom for readable text
        band_h = 2.35
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(th - band_h),
            Inches(tw),
            Inches(band_h),
        )
        band.fill.solid()
        band.fill.fore_color.rgb = tokens.color_overlay
        try:
            band.fill.transparency = 0.22
        except (AttributeError, TypeError):
            pass
        band.line.fill.background()
        text_top = th - band_h + 0.22
        safe_bottom = th - tokens.margin_in
        tb_h = max(0.85, safe_bottom - text_top - 0.06)
        tb = slide.shapes.add_textbox(
            Inches(tokens.margin_in),
            Inches(text_top),
            Inches(tw - 2 * tokens.margin_in),
            Inches(tb_h),
        )
        tf = tb.text_frame
        configure_text_frame(tf, tokens)
        clear_tf(tf)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(tokens.title_pt)
        p.font.bold = True
        p.font.name = tokens.font_title
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(tokens.subtitle_pt)
            p2.font.name = tokens.font_body
            p2.font.color.rgb = RGBColor(226, 232, 240)
            p2.space_before = Pt(10)

    else:
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(0.14),
            Inches(th),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = tokens.color_accent
        accent.line.fill.background()
        tb = slide.shapes.add_textbox(
            Inches(tokens.margin_in + 0.12),
            Inches(th * 0.28),
            Inches(tw - 2 * tokens.margin_in - 0.12),
            Inches(2.8),
        )
        tf = tb.text_frame
        configure_text_frame(tf, tokens)
        clear_tf(tf)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(tokens.title_pt)
        p.font.bold = True
        p.font.name = tokens.font_title
        p.font.color.rgb = tokens.color_text_primary
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(tokens.subtitle_pt)
        p2.font.name = tokens.font_body
        p2.font.color.rgb = tokens.color_text_muted
        p2.space_before = Pt(12)

    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes


def layout_section_divider(
    prs: Presentation,
    tokens: DesignTokens,
    *,
    section_title: str,
    kicker: str = "",
    speaker_notes: str = "",
) -> None:
    slide = add_blank_slide(prs)
    apply_slide_background(slide, tokens.color_accent_dark, tokens)
    tw = tokens.slide_width_in
    th = tokens.slide_height_in
    section_title = truncate_title(section_title, tokens.max_title_chars)
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(tokens.margin_in),
        Inches(th * 0.38),
        Inches(tw - 2 * tokens.margin_in),
        Inches(1.65),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = tokens.color_bg
    band.line.fill.background()
    tf = band.text_frame
    configure_text_frame(tf, tokens)
    clear_tf(tf)
    if kicker:
        pk = tf.paragraphs[0]
        pk.text = kicker.upper()
        pk.font.size = Pt(tokens.label_pt)
        pk.font.name = tokens.font_body
        pk.font.color.rgb = tokens.color_accent
        pk.space_after = Pt(6)
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(tokens.title_min_pt + 2)
    p.font.bold = True
    p.font.name = tokens.font_title
    p.font.color.rgb = tokens.color_text_primary
    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes


# --- 3. Image + explanation split ---
def layout_image_text_split(
    prs: Presentation,
    tokens: DesignTokens,
    tmp_dir: Path,
    temps: list[Path],
    *,
    title: str,
    bullets: Sequence[str],
    image_path: Path | None,
    image_side: str = "right",  # 'right' | 'left'
    speaker_notes: str = "",
) -> None:
    slide = add_blank_slide(prs)
    apply_slide_background(slide, tokens.color_bg, tokens)
    tw = tokens.slide_width_in
    th = tokens.slide_height_in
    m = tokens.margin_in
    title = truncate_title(title, tokens.max_title_chars)

    has_img = image_path is not None and image_path.exists()
    img_w_frac = tokens.image_panel_fraction if has_img else 0.0
    gap = tokens.column_gap_in if has_img else 0.0
    img_panel_w = tw * img_w_frac
    text_w_in = tw - 2 * m - gap - img_panel_w

    title_h = title_region_height_in(tokens, tokens.title_pt, None)
    title_box = slide.shapes.add_textbox(Inches(m), Inches(m), Inches(tw - 2 * m), Inches(title_h + 0.05))
    tf_t = title_box.text_frame
    configure_text_frame(tf_t, tokens)
    clear_tf(tf_t)
    pt = tf_t.paragraphs[0]
    pt.text = title
    pt.font.size = Pt(tokens.title_pt)
    pt.font.bold = True
    pt.font.name = tokens.font_title
    pt.font.color.rgb = tokens.color_text_primary

    body_top = m + title_h + 0.05
    body_h = th - body_top - m - 0.35

    if image_side == "right":
        text_left = m
        img_left = tw - m - img_panel_w
    else:
        img_left = m
        text_left = m + img_panel_w + tokens.column_gap_in

    work_bullets = [summarize_bullet(b, tokens.max_words_per_bullet) for b in bullets]
    fitted, font_sz, pages = fit_bullets_to_box(
        list(work_bullets),
        text_w_in,
        body_h,
        title_height_reserved_in=0,
        font_start=tokens.body_pt,
        font_min=tokens.body_min_pt,
        tokens=tokens,
    )

    if image_path and image_path.exists():
        _pic, tmp = add_image_cover(
            slide,
            image_path,
            Inches(img_left),
            Inches(body_top),
            Inches(img_panel_w),
            Inches(body_h),
            tmp_dir,
        )
        temps.append(tmp)

    body_box = slide.shapes.add_textbox(Inches(text_left), Inches(body_top), Inches(text_w_in), Inches(body_h))
    tf_b = body_box.text_frame
    configure_text_frame(tf_b, tokens)
    clear_tf(tf_b)
    for i, line in enumerate(fitted):
        p = tf_b.paragraphs[0] if i == 0 else tf_b.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(font_sz)
        p.font.name = tokens.font_body
        p.font.color.rgb = tokens.color_text_primary
        p.space_after = Pt(tokens.bullet_spacing_after_pt)

    if image_path and image_path.exists():
        cap = slide.shapes.add_textbox(Inches(m), Inches(th - m - 0.22), Inches(tw - 2 * m), Inches(0.2))
        ctf = cap.text_frame
        configure_text_frame(ctf, tokens)
        clear_tf(ctf)
        cp = ctf.paragraphs[0]
        cp.text = "Image: NASA/JPL-Caltech (cover crop — aspect preserved)"
        cp.font.size = Pt(tokens.caption_pt)
        cp.font.italic = True
        cp.font.color.rgb = tokens.color_text_muted
        cp.alignment = PP_ALIGN.RIGHT

    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes

    # Continuation slides if split required
    if len(pages) > 1:
        for chunk in pages[1:]:
            _continuation_bullets_slide(
                prs,
                tokens,
                title + " (continued)",
                chunk,
                speaker_notes="" if not speaker_notes else f"{speaker_notes} [continued]",
            )


def _continuation_bullets_slide(
    prs: Presentation,
    tokens: DesignTokens,
    title: str,
    bullets: Sequence[str],
    speaker_notes: str = "",
) -> None:
    slide = add_blank_slide(prs)
    apply_slide_background(slide, tokens.color_bg, tokens)
    tw = tokens.slide_width_in
    th = tokens.slide_height_in
    m = tokens.margin_in
    title_h = title_region_height_in(tokens, tokens.title_pt, None)
    title_box = slide.shapes.add_textbox(Inches(m), Inches(m), Inches(tw - 2 * m), Inches(title_h))
    tf_t = title_box.text_frame
    configure_text_frame(tf_t, tokens)
    clear_tf(tf_t)
    tf_t.paragraphs[0].text = truncate_title(title, tokens.max_title_chars)
    tf_t.paragraphs[0].font.size = Pt(tokens.title_pt - 2)
    tf_t.paragraphs[0].font.bold = True
    tf_t.paragraphs[0].font.name = tokens.font_title
    tf_t.paragraphs[0].font.color.rgb = tokens.color_text_primary

    body_top = m + title_h + 0.08
    body_h = th - body_top - m
    text_w = tw - 2 * m
    fitted, font_sz, _ = fit_bullets_to_box(
        list(bullets),
        text_w,
        body_h,
        0,
        tokens.body_pt,
        tokens.body_min_pt,
        tokens,
    )
    body_box = slide.shapes.add_textbox(Inches(m), Inches(body_top), Inches(text_w), Inches(body_h))
    tf_b = body_box.text_frame
    configure_text_frame(tf_b, tokens)
    clear_tf(tf_b)
    for i, line in enumerate(fitted):
        p = tf_b.paragraphs[0] if i == 0 else tf_b.add_paragraph()
        p.text = line
        p.font.size = Pt(font_sz)
        p.font.name = tokens.font_body
        p.font.color.rgb = tokens.color_text_primary
        p.space_after = Pt(tokens.bullet_spacing_after_pt)
    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes


# --- 4. Big idea ---
def layout_big_idea(
    prs: Presentation,
    tokens: DesignTokens,
    tmp_dir: Path,
    temps: list[Path],
    *,
    headline: str,
    image_path: Path | None = None,
    speaker_notes: str = "",
) -> None:
    slide = add_blank_slide(prs)
    apply_slide_background(slide, tokens.color_bg_alt, tokens)
    tw = tokens.slide_width_in
    th = tokens.slide_height_in
    m = tokens.margin_in
    headline = truncate_title(headline, 160)

    img_h = 3.15 if image_path else 0

    if image_path and image_path.exists():
        _pic, tmp = add_image_cover(
            slide,
            image_path,
            Inches(m),
            Inches(m),
            Inches(tw - 2 * m),
            Inches(img_h),
            tmp_dir,
        )
        temps.append(tmp)

    tb = slide.shapes.add_textbox(
        Inches(m),
        Inches(m + (img_h + 0.25 if image_path else 0.6)),
        Inches(tw - 2 * m),
        Inches(2.4),
    )
    tf = tb.text_frame
    configure_text_frame(tf, tokens)
    clear_tf(tf)
    p = tf.paragraphs[0]
    p.text = headline
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = tokens.font_title
    p.font.color.rgb = tokens.color_text_primary
    p.line_spacing = 1.15

    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes


# --- 5. Three-point cards ---
def layout_three_point(
    prs: Presentation,
    tokens: DesignTokens,
    *,
    title: str,
    points: Sequence[tuple[str, str]],  # (heading, blurb)
    speaker_notes: str = "",
) -> None:
    if len(points) > 3:
        points = points[:3]
    slide = add_blank_slide(prs)
    apply_slide_background(slide, tokens.color_bg, tokens)
    tw = tokens.slide_width_in
    th = tokens.slide_height_in
    m = tokens.margin_in
    title_h = title_region_height_in(tokens, tokens.title_pt, None)
    tt = slide.shapes.add_textbox(Inches(m), Inches(m), Inches(tw - 2 * m), Inches(title_h))
    tf0 = tt.text_frame
    configure_text_frame(tf0, tokens)
    clear_tf(tf0)
    tf0.paragraphs[0].text = truncate_title(title, tokens.max_title_chars)
    tf0.paragraphs[0].font.size = Pt(tokens.title_pt)
    tf0.paragraphs[0].font.bold = True
    tf0.paragraphs[0].font.name = tokens.font_title
    tf0.paragraphs[0].font.color.rgb = tokens.color_text_primary

    gap = 0.18
    card_w = (tw - 2 * m - 2 * gap) / 3
    top = m + title_h + 0.12
    card_h = th - top - m - 0.15

    for i, (hd, blurb) in enumerate(points):
        left = m + i * (card_w + gap)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(card_w),
            Inches(card_h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = tokens.color_card
        card.line.color.rgb = RGBColor(226, 232, 240)
        tf = card.text_frame
        configure_text_frame(tf, tokens)
        clear_tf(tf)
        p1 = tf.paragraphs[0]
        p1.text = hd
        p1.font.size = Pt(tokens.subtitle_pt)
        p1.font.bold = True
        p1.font.name = tokens.font_title
        p1.font.color.rgb = tokens.color_accent_dark
        p2 = tf.add_paragraph()
        p2.text = summarize_bullet(blurb, 22)
        p2.font.size = Pt(tokens.body_min_pt + 1)
        p2.font.name = tokens.font_body
        p2.font.color.rgb = tokens.color_text_primary
        p2.space_before = Pt(10)
        p2.line_spacing = 1.12

    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes


# --- 5b. Four subsystem diagram prompts (2×2 cards) ---
def layout_four_subsystem_cards(
    prs: Presentation,
    tokens: DesignTokens,
    *,
    title: str,
    subtitle: str = "",
    cards: Sequence[tuple[str, str, str]],
    speaker_notes: str = "",
) -> None:
    """
    One slide per focus: four cards for subsystem codes (e.g. W1–W4).
    Each card: `cards` item is (code, subsystem_name, diagram_must_show_summary).
    """
    if len(cards) != 4:
        raise ValueError("layout_four_subsystem_cards requires exactly 4 cards")

    slide = add_blank_slide(prs)
    apply_slide_background(slide, tokens.color_bg, tokens)
    tw = tokens.slide_width_in
    th = tokens.slide_height_in
    m = tokens.margin_in

    title_h = title_region_height_in(tokens, tokens.title_pt, None)
    tt = slide.shapes.add_textbox(Inches(m), Inches(m), Inches(tw - 2 * m), Inches(title_h + (0.35 if subtitle else 0)))
    tf0 = tt.text_frame
    configure_text_frame(tf0, tokens)
    clear_tf(tf0)
    tf0.paragraphs[0].text = truncate_title(title, tokens.max_title_chars)
    tf0.paragraphs[0].font.size = Pt(tokens.title_pt)
    tf0.paragraphs[0].font.bold = True
    tf0.paragraphs[0].font.name = tokens.font_title
    tf0.paragraphs[0].font.color.rgb = tokens.color_text_primary
    if subtitle:
        sp = tf0.add_paragraph()
        sp.text = subtitle
        sp.font.size = Pt(tokens.subtitle_min_pt)
        sp.font.name = tokens.font_body
        sp.font.color.rgb = tokens.color_text_muted
        sp.space_before = Pt(6)

    top = m + title_h + (0.42 if subtitle else 0.18)
    gap_x = 0.2
    gap_y = 0.16
    grid_w = tw - 2 * m
    grid_h = th - top - m - 0.1
    card_w = (grid_w - gap_x) / 2
    card_h = (grid_h - gap_y) / 2

    positions = [
        (m, top),
        (m + card_w + gap_x, top),
        (m, top + card_h + gap_y),
        (m + card_w + gap_x, top + card_h + gap_y),
    ]

    for (left, card_top), (code, sub_name, must_show) in zip(positions, cards):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(card_top),
            Inches(card_w),
            Inches(card_h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = tokens.color_card
        card.line.color.rgb = RGBColor(226, 232, 240)
        tf = card.text_frame
        configure_text_frame(tf, tokens)
        clear_tf(tf)
        h1 = tf.paragraphs[0]
        h1.text = f"{code} — {sub_name}"
        h1.font.size = Pt(tokens.label_pt + 1)
        h1.font.bold = True
        h1.font.name = tokens.font_title
        h1.font.color.rgb = tokens.color_accent_dark
        h2 = tf.add_paragraph()
        h2.text = summarize_bullet(must_show, 36)
        h2.font.size = Pt(tokens.body_min_pt)
        h2.font.name = tokens.font_body
        h2.font.color.rgb = tokens.color_text_primary
        h2.space_before = Pt(6)
        h2.line_spacing = 1.12

    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes


# --- 6. Process / timeline ---
def layout_process_timeline(
    prs: Presentation,
    tokens: DesignTokens,
    *,
    title: str,
    steps: Sequence[str],
    speaker_notes: str = "",
) -> None:
    steps = list(steps)[:5]
    slide = add_blank_slide(prs)
    apply_slide_background(slide, tokens.color_bg, tokens)
    tw = tokens.slide_width_in
    th = tokens.slide_height_in
    m = tokens.margin_in
    title_h = title_region_height_in(tokens, tokens.title_pt, None)
    tt = slide.shapes.add_textbox(Inches(m), Inches(m), Inches(tw - 2 * m), Inches(title_h))
    tf0 = tt.text_frame
    configure_text_frame(tf0, tokens)
    clear_tf(tf0)
    tf0.paragraphs[0].text = truncate_title(title, tokens.max_title_chars)
    tf0.paragraphs[0].font.size = Pt(tokens.title_pt)
    tf0.paragraphs[0].font.bold = True
    tf0.paragraphs[0].font.name = tokens.font_title
    tf0.paragraphs[0].font.color.rgb = tokens.color_text_primary

    y = m + title_h + 0.15
    for idx, st in enumerate(steps, start=1):
        row_h = 0.85
        num = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(m), Inches(y), Inches(0.42), Inches(0.42))
        num.fill.solid()
        num.fill.fore_color.rgb = tokens.color_accent
        num.line.fill.background()
        ntf = num.text_frame
        configure_text_frame(ntf, tokens)
        clear_tf(ntf)
        np = ntf.paragraphs[0]
        np.text = str(idx)
        np.font.size = Pt(16)
        np.font.bold = True
        np.font.color.rgb = RGBColor(255, 255, 255)
        np.alignment = PP_ALIGN.CENTER
        tb = slide.shapes.add_textbox(Inches(m + 0.55), Inches(y - 0.02), Inches(tw - 2 * m - 0.55), Inches(row_h))
        tf = tb.text_frame
        configure_text_frame(tf, tokens)
        clear_tf(tf)
        p = tf.paragraphs[0]
        p.text = summarize_bullet(st, 24)
        p.font.size = Pt(tokens.body_pt)
        p.font.name = tokens.font_body
        p.font.color.rgb = tokens.color_text_primary
        y += row_h + 0.12

    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes


# --- 7. Quote ---
def layout_quote(
    prs: Presentation,
    tokens: DesignTokens,
    *,
    quote: str,
    source: str,
    speaker_notes: str = "",
) -> None:
    slide = add_blank_slide(prs)
    apply_slide_background(slide, tokens.color_bg_alt, tokens)
    tw = tokens.slide_width_in
    th = tokens.slide_height_in
    m = tokens.margin_in
    quote = summarize_bullet(quote, 48) if len(quote.split()) > 48 else quote
    box = slide.shapes.add_textbox(Inches(m + 0.2), Inches(th * 0.22), Inches(tw - 2 * m - 0.4), Inches(3.2))
    tf = box.text_frame
    configure_text_frame(tf, tokens)
    clear_tf(tf)
    p = tf.paragraphs[0]
    p.text = f"“{quote.strip()}”"
    p.font.size = Pt(26)
    p.font.italic = True
    p.font.name = tokens.font_body
    p.font.color.rgb = tokens.color_text_primary
    p.line_spacing = 1.22

    src = slide.shapes.add_textbox(Inches(m), Inches(th * 0.62), Inches(tw - 2 * m), Inches(0.9))
    stf = src.text_frame
    configure_text_frame(stf, tokens)
    clear_tf(stf)
    sp = stf.paragraphs[0]
    sp.text = source
    sp.font.size = Pt(tokens.caption_pt + 2)
    sp.font.name = tokens.font_body
    sp.font.color.rgb = tokens.color_text_muted

    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes


# --- 8. Comparison ---
def layout_comparison(
    prs: Presentation,
    tokens: DesignTokens,
    *,
    title: str,
    left_title: str,
    left_bullets: Sequence[str],
    right_title: str,
    right_bullets: Sequence[str],
    speaker_notes: str = "",
) -> None:
    slide = add_blank_slide(prs)
    apply_slide_background(slide, tokens.color_bg, tokens)
    tw = tokens.slide_width_in
    th = tokens.slide_height_in
    m = tokens.margin_in
    title_h = title_region_height_in(tokens, tokens.title_pt, None)
    tt = slide.shapes.add_textbox(Inches(m), Inches(m), Inches(tw - 2 * m), Inches(title_h))
    tf0 = tt.text_frame
    configure_text_frame(tf0, tokens)
    clear_tf(tf0)
    tf0.paragraphs[0].text = truncate_title(title, tokens.max_title_chars)
    tf0.paragraphs[0].font.size = Pt(tokens.title_pt)
    tf0.paragraphs[0].font.bold = True
    tf0.paragraphs[0].font.name = tokens.font_title

    col_w = (tw - 2 * m - tokens.column_gap_in) / 2
    top = m + title_h + 0.12
    col_h = th - top - m

    def add_col(left_in: float, label: str, items: Sequence[str]):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_in),
            Inches(top),
            Inches(col_w),
            Inches(0.38),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = tokens.color_accent
        bar.line.fill.background()
        btf = bar.text_frame
        configure_text_frame(btf, tokens)
        clear_tf(btf)
        bp = btf.paragraphs[0]
        bp.text = label
        bp.font.size = Pt(tokens.subtitle_pt - 2)
        bp.font.bold = True
        bp.font.color.rgb = RGBColor(255, 255, 255)
        bp.font.name = tokens.font_body

        body = slide.shapes.add_textbox(Inches(left_in), Inches(top + 0.45), Inches(col_w), Inches(col_h - 0.5))
        tf = body.text_frame
        configure_text_frame(tf, tokens)
        clear_tf(tf)
        fits = [summarize_bullet(x, tokens.max_words_per_bullet) for x in items][:8]
        fs = Pt(12) if len(fits) > 5 else Pt(tokens.body_min_pt + 1)
        sp_after = Pt(4) if len(fits) > 5 else Pt(8)
        for i, line in enumerate(fits):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = fs
            p.font.name = tokens.font_body
            p.space_after = sp_after

    add_col(m, left_title, left_bullets)
    add_col(m + col_w + tokens.column_gap_in, right_title, right_bullets)

    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes


def init_widescreen_deck(tokens: DesignTokens = DEFAULT_TOKENS) -> Presentation:
    prs = Presentation()
    prs.slide_width = tokens.slide_width()
    prs.slide_height = tokens.slide_height()
    return prs

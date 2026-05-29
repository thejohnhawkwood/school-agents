"""Extract teacher Digestion & Respiratory unit inputs to run folder."""
from __future__ import annotations

from pathlib import Path

from docx import Document
from pptx import Presentation

OUT = (
    Path(__file__).resolve().parents[1]
    / "outputs"
    / "tests"
    / "bio30-digest-resp-mc30-v1"
    / "_extracted"
)

PATHS = [
    Path(r"c:\Users\pbird\Downloads\Digestion and Respiration Unit Test.docx"),
    Path(
        r"c:\Users\pbird\Downloads\1. Bio 20\Unit B-Digestion and Respiratory System\1 - Digestion.pptx"
    ),
    Path(
        r"c:\Users\pbird\Downloads\1. Bio 20\Unit B-Digestion and Respiratory System\2 - Respiratory System [Autosaved].pptx"
    ),
]


def dump_docx(p: Path) -> None:
    d = Document(str(p))
    lines = [para.text for para in d.paragraphs if para.text.strip()]
    for tbl in d.tables:
        for row in tbl.rows:
            lines.append(" | ".join(c.text.strip() for c in row.cells))
    (OUT / f"{p.stem}.txt").write_text("\n".join(lines), encoding="utf-8")


def dump_pptx(p: Path) -> None:
    prs = Presentation(str(p))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides):
        lines.append(f"\n--- slide {i + 1} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    (OUT / f"{p.stem}.txt").write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    for p in PATHS:
        if not p.exists():
            print("MISSING", p)
            continue
        print("OK", p.name)
        if p.suffix.lower() == ".docx":
            dump_docx(p)
        elif p.suffix.lower() == ".pptx":
            dump_pptx(p)
    print("wrote under", OUT)


if __name__ == "__main__":
    main()

"""Extract text from the Unit 3 review docx and the Unit C pptx for review-set planning."""
from pathlib import Path
from docx import Document
from pptx import Presentation

DOWNLOADS = Path(r"C:\Users\pbird\Downloads")
DOCX = DOWNLOADS / "BONUS - Unit 3 Review Sheet.docx"
# Resolve pptx dynamically because the filename contains a unicode dash
PPTX_CANDIDATES = [
    f for f in DOWNLOADS.iterdir()
    if f.suffix == ".pptx"
    and "Electromagnetic" in f.name
    and not f.name.startswith("~$")
]
assert PPTX_CANDIDATES, "Could not find Unit C pptx in Downloads"
PPTX = PPTX_CANDIDATES[0]

OUT = Path(__file__).resolve().parent.parent / "outputs" / "_extracted"
OUT.mkdir(parents=True, exist_ok=True)


def dump_docx(path: Path, out_file: Path) -> None:
    doc = Document(path)
    lines: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            lines.append(text)
    for ti, table in enumerate(doc.tables):
        lines.append(f"\n[TABLE {ti}]")
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " | ") for c in row.cells]
            lines.append(" || ".join(cells))
    out_file.write_text("\n".join(lines), encoding="utf-8")
    print(f"Wrote {out_file} ({len(lines)} lines)")


def dump_pptx(path: Path, out_file: Path) -> None:
    prs = Presentation(path)
    lines: list[str] = []
    for si, slide in enumerate(prs.slides, start=1):
        lines.append(f"\n===== SLIDE {si} =====")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        lines.append(text)
            if shape.shape_type == 19:  # table
                try:
                    tbl = shape.table
                    for row in tbl.rows:
                        cells = [c.text.strip().replace("\n", " | ") for c in row.cells]
                        lines.append(" || ".join(cells))
                except Exception:
                    pass
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"[NOTES] {notes}")
    out_file.write_text("\n".join(lines), encoding="utf-8")
    print(f"Wrote {out_file} ({len(lines)} lines)")


if __name__ == "__main__":
    dump_docx(DOCX, OUT / "unit3_review_sheet.txt")
    dump_pptx(PPTX, OUT / "unit_c_electromagnetic_energy.txt")

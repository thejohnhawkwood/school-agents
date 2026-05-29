"""One-off: extract text from teacher-provided PS&CR inputs into outputs/tests."""
from __future__ import annotations

from pathlib import Path

from docx import Document
from pypdf import PdfReader
from pptx import Presentation

OUT = Path(__file__).resolve().parents[1] / "outputs" / "tests" / "bio30-pscr-mc30-v1" / "_extracted"

PATHS = [
    Path(r"c:\Users\pbird\Downloads\PS&CR Unit A Exam.docx"),
    Path(r"c:\Users\pbird\Downloads\PS&CR Unit C Exam with answers.pdf"),
    Path(r"c:\Users\pbird\Downloads\pos-bio-20-30 (2).pdf"),
    Path(r"c:\Users\pbird\Downloads\Photosynthesis & Cell Respiration Bird 2026.pptx"),
]


def dump_docx(p: Path) -> None:
    d = Document(str(p))
    lines = [para.text for para in d.paragraphs if para.text.strip()]
    for tbl in d.tables:
        for row in tbl.rows:
            lines.append(" | ".join(c.text.strip() for c in row.cells))
    (OUT / f"{p.stem}.txt").write_text("\n".join(lines), encoding="utf-8")


def dump_pdf(p: Path) -> None:
    r = PdfReader(str(p))
    parts: list[str] = []
    for i, page in enumerate(r.pages):
        t = page.extract_text() or ""
        parts.append(f"\n--- page {i + 1} ---\n{t}")
    (OUT / f"{p.stem}.txt").write_text("".join(parts), encoding="utf-8")


def dump_pptx(p: Path) -> None:
    prs = Presentation(str(p))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides):
        lines.append(f"\n--- slide {i + 1} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    (OUT / f"{p.stem}.txt").write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    for p in PATHS:
        if not p.exists():
            print("MISSING", p)
            continue
        print("OK", p.name)
        suf = p.suffix.lower()
        if suf == ".docx":
            dump_docx(p)
        elif suf == ".pdf":
            dump_pdf(p)
        elif suf == ".pptx":
            dump_pptx(p)
    print("wrote under", OUT)


if __name__ == "__main__":
    main()

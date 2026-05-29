"""Slide QA — bounding boxes, minimum font, safe area."""

from __future__ import annotations

from dataclasses import dataclass, field

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx_kit.tokens import DesignTokens, DEFAULT_TOKENS


@dataclass
class SlideQAReport:
    slide_index: int
    errors: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)

    @property
    def ok(self) -> bool:
        return len(self.errors) == 0


def _collect_shapes(shapes):
    """Flatten group shapes."""
    out = []
    for sh in shapes:
        out.append(sh)
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            out.extend(_collect_shapes(sh.shapes))
    return out


def _shape_text_stripped(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    try:
        parts = []
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                parts.append(r.text)
            if not p.runs:
                parts.append(p.text)
        return "".join(parts).strip()
    except AttributeError:
        return ""


def assertWithinSafeArea(
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    tokens: DesignTokens = DEFAULT_TOKENS,
) -> None:
    """Raises ValueError if axis-aligned box intrudes past slide margins."""
    m = tokens.margin_in
    sw, sh = tokens.slide_width_in, tokens.slide_height_in
    if (
        left_in < m - 1e-6
        or top_in < m - 1e-6
        or left_in + width_in > sw - m + 1e-6
        or top_in + height_in > sh - m + 1e-6
    ):
        raise ValueError(
            f"Box ({left_in:.2f},{top_in:.2f},{width_in:.2f},{height_in:.2f}) outside safe margin {m}"
        )


def validate_slide(
    slide,
    slide_index: int,
    tokens: DesignTokens = DEFAULT_TOKENS,
    eps_in: float = 0.02,
) -> SlideQAReport:
    report = SlideQAReport(slide_index=slide_index)
    sw = tokens.slide_width_in + eps_in
    sh = tokens.slide_height_in + eps_in
    margin = tokens.margin_in

    for shape in _collect_shapes(slide.shapes):
        try:
            l = shape.left.inches
            t = shape.top.inches
            w = shape.width.inches
            h = shape.height.inches
        except AttributeError:
            continue

        if l + w > sw:
            report.errors.append(f"Shape extends past right edge: left={l:.2f} w={w:.2f}")
        if t + h > sh:
            report.errors.append(f"Shape extends past bottom edge: top={t:.2f} h={h:.2f}")
        if l < -eps_in or t < -eps_in:
            report.errors.append(f"Shape negative origin: left={l:.2f} top={t:.2f}")

        # Safe area (full-bleed images may intentionally touch edges — warn only if text box)
        if getattr(shape, "has_text_frame", False) and _shape_text_stripped(shape):
            if l < margin - eps_in or t < margin - eps_in:
                report.warnings.append(
                    f"Text box near/outside safe margin (left/top): ({l:.2f},{t:.2f})"
                )
            if l + w > tokens.slide_width_in - margin + eps_in:
                report.warnings.append("Text box extends into right margin")
            if t + h > tokens.slide_height_in - margin + eps_in:
                report.warnings.append("Text box extends into bottom margin")

        if getattr(shape, "has_text_frame", False):
            try:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.size:
                            pt = run.font.size.pt
                            if pt < 10:
                                report.errors.append(f"Font below absolute minimum: {pt:.1f} pt")
                            elif pt < tokens.body_min_pt - 0.5:
                                report.warnings.append(
                                    f"Font below preferred body minimum ({tokens.body_min_pt} pt): {pt:.1f} pt"
                                )
            except AttributeError:
                pass

    return report


def validate_deck(prs, tokens: DesignTokens = DEFAULT_TOKENS) -> list[SlideQAReport]:
    reports = []
    for i, slide in enumerate(prs.slides):
        reports.append(validate_slide(slide, i, tokens))
    return reports
